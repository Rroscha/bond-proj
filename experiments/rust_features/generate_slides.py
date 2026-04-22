#!/usr/bin/env python3
"""Generate a .pptx slide deck for the RustDiff project presentation (~7 min)."""
import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE

DATA = json.loads(Path('experiments/rust_features/results/analysis_data.json').read_text())

# ── Colors (white background theme) ──
BG       = RGBColor(0xFF, 0xFF, 0xFF)
WHITE    = RGBColor(0x1A, 0x1A, 0x2E)  # dark text on white bg
BLUE     = RGBColor(0x1A, 0x6B, 0xD4)
PURPLE   = RGBColor(0x7C, 0x3A, 0xED)
GREEN    = RGBColor(0x16, 0x7F, 0x39)
RED      = RGBColor(0xCF, 0x22, 0x2E)
RUST_CLR = RGBColor(0xC4, 0x34, 0x3B)
C_CLR    = RGBColor(0x0E, 0x5A, 0xAB)
GRAY     = RGBColor(0x65, 0x6D, 0x76)
DARK_BG  = RGBColor(0xF5, 0xF5, 0xF5)  # light gray for code boxes
YELLOW   = RGBColor(0xB4, 0x7A, 0x00)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


def set_slide_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_frame(slide, left, top, width, height, items, font_size=16,
                     color=WHITE, bold_first=False, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = spacing
        if isinstance(item, tuple):
            txt, clr = item
        else:
            txt, clr = item, color
        p.text = txt
        p.font.size = Pt(font_size)
        p.font.color.rgb = clr
        p.font.name = 'Calibri'
        p.level = 0
        if bold_first and i == 0:
            p.font.bold = True
    return tf


def add_code_box(slide, left, top, width, height, code, font_size=11, title=None, title_color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BG
    shape.line.color.rgb = RGBColor(0xD0, 0xD7, 0xDE)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(6)
    if title:
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.color.rgb = title_color
        p.font.bold = True
        p.font.name = 'Calibri'
        p.space_after = Pt(4)
        p2 = tf.add_paragraph()
    else:
        p2 = tf.paragraphs[0]
    p2.text = code
    p2.font.size = Pt(font_size)
    p2.font.color.rgb = RGBColor(0x24, 0x29, 0x2E)  # dark code text
    p2.font.name = 'Consolas'
    return shape


def add_rect(slide, left, top, width, height, fill_color, text='', font_size=12, font_color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(6)
        tf.margin_top = Pt(2)
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.name = 'Calibri'
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    return shape


def slide_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.8),
                title, font_size=32, color=BLUE, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.6), Inches(1.0), Inches(12), Inches(0.5),
                    subtitle, font_size=18, color=GRAY)


# ═══════════════════════════════════════════════════════════════════════════
# Slide 1: Title
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
add_textbox(sl, Inches(1), Inches(2.0), Inches(11), Inches(1.2),
            'RustDiff', font_size=54, color=BLUE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(sl, Inches(1), Inches(3.2), Inches(11), Inches(0.8),
            'How Rust Language Features Affect Binary Block Matching',
            font_size=26, color=WHITE, alignment=PP_ALIGN.CENTER)
add_textbox(sl, Inches(1), Inches(4.2), Inches(11), Inches(0.6),
            'Comparing 3 matching features across Rust vs. C  |  100 function pairs  |  DWARF ground truth',
            font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)
add_textbox(sl, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
            'Rroscha  •  Spring 2026',
            font_size=18, color=GRAY, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# Slide 2: Research Question
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
slide_title(sl, 'Research Question')

add_textbox(sl, Inches(0.8), Inches(1.6), Inches(11.5), Inches(1.0),
            'Which block-level matching features remain stable across\n'
            'compiler optimizations (O0 → O2), and how does Rust\'s\n'
            'safety overhead affect each feature compared to equivalent C?',
            font_size=24, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_bullet_frame(sl, Inches(0.8), Inches(3.6), Inches(5.5), Inches(3.5), [
    ('Why this matters:', BLUE),
    ('• Binary diffing is key for security analysis, patch diffing, malware comparison', WHITE),
    ('• Tools like BinDiff & Diaphora were designed with C/C++ in mind', WHITE),
    ('• Rust adoption is growing fast in systems software', WHITE),
    ('• Do existing matching features still work on Rust binaries?', YELLOW),
], font_size=16)

add_bullet_frame(sl, Inches(7.0), Inches(3.6), Inches(5.5), Inches(3.5), [
    ('Approach:', BLUE),
    ('• 100 Rust functions + 100 equivalent C functions', WHITE),
    ('• Same algorithm, different language idioms', WHITE),
    ('• Compile each at O0 and O2 with debug info', WHITE),
    ('• Match O0↔O2 blocks, validate with DWARF ground truth', WHITE),
], font_size=16)


# ═══════════════════════════════════════════════════════════════════════════
# Slide 3: Experiment Setup
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
slide_title(sl, 'Experiment Setup')

# Left: 5 features
add_textbox(sl, Inches(0.6), Inches(1.5), Inches(5.5), Inches(0.5),
            '5 Rust-specific features tested (20 functions each):',
            font_size=16, color=BLUE, bold=True)

features_data = [
    ('Bounds Checking (bc_)', 'data[i] → cmp + jae + panic'),
    ('Safe Abstractions (om_)', 'ownership model → Vec, Iterator, Drop generate extra blocks'),
    ('Drop Glue / RAII (dg_)', 'auto-generated drop_in_place<T> calls'),
    ('? Operator (qm_)', 'Result/Option → discriminant check + early return'),
    ('Panic Paths (pu_)', '.unwrap() → hidden panic call blocks'),
]
y = 2.2
for feat, note in features_data:
    add_textbox(sl, Inches(0.8), Inches(y), Inches(5.0), Inches(0.35),
                f'• {feat}', font_size=15, color=WHITE, bold=True)
    add_textbox(sl, Inches(1.0), Inches(y + 0.3), Inches(4.8), Inches(0.3),
                note, font_size=13, color=GRAY)
    y += 0.65

# Right: 3 matching features
add_textbox(sl, Inches(6.8), Inches(1.5), Inches(6.0), Inches(0.5),
            '3 matching features compared:', font_size=16, color=BLUE, bold=True)

add_textbox(sl, Inches(6.8), Inches(2.0), Inches(6.0), Inches(0.4),
            'These are standard features used in binary diffing\n'
            '(BinDiff, Diaphora, etc.) — not novel to this work.',
            font_size=13, color=GRAY)

methods = [
    ('Value-based (micro-execution)', 'Run each block with test inputs via angr;\ncompare register outputs, dataflow, memory pattern, constants'),
    ('Opcode similarity (Jaccard)', 'Compare instruction type sets:\noverlap / union of {mov, cmp, add, ...}'),
    ('Constant similarity (Jaccard)', 'Compare embedded numeric constants:\noverlap / union of immediates'),
]
y = 2.7
for title, desc in methods:
    add_textbox(sl, Inches(7.0), Inches(y), Inches(5.5), Inches(0.35),
                f'• {title}', font_size=15, color=WHITE, bold=True)
    add_textbox(sl, Inches(7.2), Inches(y + 0.3), Inches(5.3), Inches(0.5),
                desc, font_size=13, color=GRAY)
    y += 0.75

# Bottom: Ground truth
add_rect(sl, Inches(0.6), Inches(5.8), Inches(12.1), Inches(1.2), DARK_BG,
         '', font_size=14)
add_textbox(sl, Inches(0.8), Inches(5.9), Inches(11.7), Inches(0.4),
            'Ground Truth: DWARF debug info', font_size=16, color=GREEN, bold=True)
add_textbox(sl, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.5),
            'O0 and O2 binaries both built with DWARF symbols enabled (does not affect optimization).\n'
            'Each block is mapped to source lines. Two blocks are "correctly matched" iff they share ≥1 source line.\n'
            'Accuracy = correct pairs / total pairs (Hungarian algorithm, one-to-one).',
            font_size=14, color=WHITE)

# ═══════════════════════════════════════════════════════════════════════════
# Helper: extract real asm from data (needed by Slide 4 onward)
# ═══════════════════════════════════════════════════════════════════════════
def get_fn(name, lang):
    for f in DATA['functions']:
        if f['base_name'] == name and f['lang'] == lang:
            return f
    return {}

def get_blocks(fn_data, side):
    return [b for b in fn_data.get('blocks', []) if b['side'] == side]

def get_blocks_by_type(fn_data, side, btype):
    return [b for b in fn_data.get('blocks', []) if b['side'] == side and btype in b['type']]

def fmt_asm(block, max_lines=6):
    lines = []
    for line in block.get('asm', [])[:max_lines]:
        parts = line.split(': ', 1)
        lines.append(parts[-1] if len(parts) > 1 else line)
    if len(block.get('asm', [])) > max_lines:
        lines.append(f'... +{len(block["asm"]) - max_lines} more')
    return '\n'.join(lines)

def fmt_asm_raw(block, max_lines=5):
    lines = block.get('asm', [])[:max_lines]
    out = '\n'.join(lines)
    if len(block.get('asm', [])) > max_lines:
        out += f'\n... +{len(block["asm"]) - max_lines} more'
    return out

def get_pair_example(fn_data, method, want_correct):
    blocks = {b['addr']: b for b in fn_data.get('blocks', [])}
    pairs = fn_data['methods'][method]['pairs']
    for p in pairs:
        if p['is_correct'] == want_correct:
            b0 = blocks.get(p['addr_o0'])
            b2 = blocks.get(p['addr_o2'])
            if b0 and b2 and b0.get('asm') and b2.get('asm'):
                return b0, b2, p
    return None, None, None


# ═══════════════════════════════════════════════════════════════════════════
# Slide 4: How block matching works (diagram-style)
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
slide_title(sl, 'How It Works', 'From source code to matched blocks')

# Pipeline: Source → Compile O0/O2 → Extract blocks → Compute features → Hungarian → Validate
steps = [
    ('Source\nCode', PURPLE),
    ('Compile\nO0 & O2', BLUE),
    ('Extract\nBlocks', C_CLR),
    ('Compute\nFeatures', GREEN),
    ('Hungarian\nMatching', YELLOW),
    ('DWARF\nValidation', RUST_CLR),
]
box_w = Inches(1.7)
box_h = Inches(1.0)
start_x = Inches(0.5)
y_pos = Inches(2.0)
gap = Inches(0.3)
for i, (label, color) in enumerate(steps):
    x = start_x + i * (box_w + gap)
    add_rect(sl, x, y_pos, box_w, box_h, color, label, font_size=14,
             font_color=RGBColor(0xFF, 0xFF, 0xFF))
    if i < len(steps) - 1:
        arrow_x = x + box_w + Pt(2)
        add_textbox(sl, arrow_x, y_pos + Inches(0.2), gap, Inches(0.6),
                    '→', font_size=28, color=GRAY, alignment=PP_ALIGN.CENTER)

# Example: what features look like
add_textbox(sl, Inches(0.6), Inches(3.5), Inches(12), Inches(0.5),
            'Feature extraction example — one block:', font_size=16, color=BLUE, bold=True)

# Use a real BOUNDS_CHECK block from bc_01 Rust O0
_bc_r = get_fn('bc_01', 'Rust')
_bc_blk = get_blocks_by_type(_bc_r, 'O0', 'BOUNDS_CHECK')
_real_blk = _bc_blk[0] if _bc_blk else get_blocks(_bc_r, 'O0')[0]

add_code_box(sl, Inches(0.6), Inches(4.0), Inches(3.8), Inches(2.8),
             fmt_asm_raw(_real_blk, 6),
             font_size=12, title=f'Real block: bc_01 Rust O0 — {_real_blk["type"]} ({_real_blk["n_insns"]} insns)',
             title_color=BLUE)

_opcodes = set()
for line in _real_blk.get('asm', []):
    parts = line.split(': ', 1)
    instr = parts[-1].strip().split()[0] if parts else ''
    if instr:
        _opcodes.add(instr)
_consts = _real_blk.get('constants', [])

add_code_box(sl, Inches(4.8), Inches(4.0), Inches(3.8), Inches(2.8),
             f'Opcodes:    {{{", ".join(sorted(_opcodes))}}}\n'
             f'Constants:  {{{", ".join(str(c) for c in _consts[:6])}}}\n\n'
             f'Value-based: run with 4 test inputs,\n'
             f'  collect register outputs + memory\n'
             f'  access patterns + dataflow edges,\n'
             f'  weighted combination (Jaccard)',
             font_size=12, title='Extracted Features', title_color=GREEN)

add_code_box(sl, Inches(9.0), Inches(4.0), Inches(3.8), Inches(2.8),
             'For each (O0_block, O2_block):\n'
             '  sim = Jaccard(feature_O0, feature_O2)\n\n'
             'Build cost matrix → Hungarian\n'
             '→ optimal 1-to-1 assignment\n\n'
             'Check each pair against DWARF\n'
             '→ accuracy = correct / total',
             font_size=13, title='Matching', title_color=YELLOW)

# ═══════════════════════════════════════════════════════════════════════════
# Slide 5: The core problem — block count explosion
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
slide_title(sl, 'The Core Problem: Safety Blocks Look Too Similar',
            'Rust\'s safety features generate blocks that are hard to tell apart')

rust_fns = [f for f in DATA['functions'] if f['lang'] == 'Rust']
c_fns = [f for f in DATA['functions'] if f['lang'] == 'C']
r_avg_o0 = sum(f['n_o0'] for f in rust_fns) / len(rust_fns)
c_avg_o0 = sum(f['n_o0'] for f in c_fns) / len(c_fns)
r_avg_o2 = sum(f['n_o2'] for f in rust_fns) / len(rust_fns)
c_avg_o2 = sum(f['n_o2'] for f in c_fns) / len(c_fns)

# Left: the problem
add_textbox(sl, Inches(0.6), Inches(1.8), Inches(5.8), Inches(0.4),
            'Problem 1: Safety blocks look alike', font_size=18, color=YELLOW, bold=True)

safety_blocks = [
    ('Bounds check', 'cmp rax, rsi; jae panic', 'All look like: cmp + conditional jump', RUST_CLR),
    ('Drop glue', 'lea rdi, [rsp+X]; call drop_in_place', 'All look like: lea + call', YELLOW),
    ('Panic path', 'lea rdi, [panic_msg]; call panic', 'All look like: lea + call', RED),
    ('Discriminant check', 'test rax, rax; je err_path', 'All look like: test + conditional jump', PURPLE),
]
y = 2.4
for title, asm, problem, color in safety_blocks:
    add_textbox(sl, Inches(0.8), Inches(y), Inches(2.2), Inches(0.3),
                f'• {title}', font_size=14, color=color, bold=True)
    add_textbox(sl, Inches(3.0), Inches(y), Inches(3.3), Inches(0.3),
                problem, font_size=13, color=GRAY)
    y += 0.4

add_textbox(sl, Inches(0.6), Inches(y + 0.2), Inches(5.8), Inches(0.5),
            'These blocks also look like computation blocks — same opcodes\n'
            '(mov, cmp, call, lea). The matcher can\'t tell them apart.',
            font_size=14, color=WHITE)

# Right: high frequency
add_textbox(sl, Inches(7.0), Inches(1.8), Inches(5.8), Inches(0.4),
            'Problem 2: They appear everywhere', font_size=18, color=YELLOW, bold=True)

add_bullet_frame(sl, Inches(7.0), Inches(2.4), Inches(5.8), Inches(2.0), [
    ('Every data[i] → a bounds check block', RUST_CLR),
    ('Every scope exit → drop glue blocks', YELLOW),
    ('Every .unwrap() → a panic block', RED),
    ('Every ? operator → a discriminant check', PURPLE),
], font_size=15, spacing=Pt(10))

add_textbox(sl, Inches(7.0), Inches(4.2), Inches(5.8), Inches(0.5),
            'A single function can trigger dozens of these.\n'
            'Each one adds another near-identical block to match.',
            font_size=14, color=WHITE)

# Bottom: consequence
add_rect(sl, Inches(0.6), Inches(5.2), Inches(12.1), Inches(2.0), DARK_BG)
add_textbox(sl, Inches(0.8), Inches(5.3), Inches(11.7), Inches(0.4),
            'Consequence: wrong matches', font_size=16, color=RED, bold=True)
add_bullet_frame(sl, Inches(0.8), Inches(5.8), Inches(11.7), Inches(1.2), [
    ('1. Many blocks with near-identical features → algorithm pairs the wrong ones with high confidence', WHITE),
    ('2. After O2 optimization, some safety blocks get merged/eliminated → the O0 block has no true partner, '
     'but the algorithm still assigns it to something that looks similar', WHITE),
    ('3. C doesn\'t have these — its blocks are more diverse, so features can distinguish them', WHITE),
], font_size=14, spacing=Pt(4))


# ═══════════════════════════════════════════════════════════════════════════
# Helper: show COMPLETE assembly for a block (no truncation)
# ═══════════════════════════════════════════════════════════════════════════
def fmt_asm_full(block):
    lines = []
    for line in block.get('asm', []):
        parts = line.split(': ', 1)
        lines.append(parts[-1] if len(parts) > 1 else line)
    return '\n'.join(lines)

def asm_box_height(block):
    n = len(block.get('asm', []))
    return Inches(0.28 + n * 0.16)

# Layout: Left = Rust asm | Right = C asm | Far right = results
LEFT   = Inches(0.3)
MID    = Inches(6.3)
BLK_W  = Inches(5.8)
RES_W  = Inches(3.1)



# ═══════════════════════════════════════════════════════════════════════════
# Slide 6: Example 1 — Bounds Checking (bc_08)
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
slide_title(sl, 'Example: Bounds Checking (bc_08)',
            'Every data[x] generates an identical bounds check block — C has none')

bc_r = get_fn('bc_08', 'Rust')
bc_c = get_fn('bc_08', 'C')

# ── Left: Rust source + assembly ──
add_code_box(sl, LEFT, Inches(1.5), BLK_W, Inches(1.8),
    'for i in 1..n-1 {\n'
    '    if data[i] > data[i-1]     // data[i]: cmp i,len; jae panic\n'
    '    && data[i] > data[i+1] {   // data[i-1]: cmp (i-1),len; jae panic\n'
    '        let mut lmin = data[i]; // data[i+1]: cmp (i+1),len; jae panic\n'
    '        for j in (0..i).rev() {\n'
    '            if data[j] < lmin   // data[j]: cmp j,len; jae panic\n'
    '// Each data[x] → compiler inserts: load index, load len,\n'
    '//   cmp index,len; jae panic_handler. Success → access memory.',
    font_size=10, title='Rust source — every data[x] inserts a bounds check (index < len)', title_color=RUST_CLR)

r_bc_o0 = get_blocks_by_type(bc_r, 'O0', 'BOUNDS_CHECK')
_y = Inches(3.15)
add_textbox(sl, LEFT, _y, BLK_W, Inches(0.3),
            'Rust assembly — 3 of 9 bounds check blocks (cmp index,len; jae panic):', font_size=12, color=YELLOW, bold=True)
_y += Inches(0.3)
bc_labels = ['← data[i] access: cmp index, len', '← data[i-1] access: cmp index, len', '← data[i+1] access: cmp index, len']
for idx in range(min(3, len(r_bc_o0))):
    blk = r_bc_o0[idx]
    h = asm_box_height(blk)
    add_code_box(sl, LEFT, _y, Inches(5.0), h,
        fmt_asm_full(blk), font_size=9,
        title=bc_labels[idx] if idx < len(bc_labels) else '',
        title_color=RUST_CLR)
    _y += h + Inches(0.05)

add_textbox(sl, LEFT, _y, BLK_W, Inches(0.5),
            '↑ All 9 blocks: identical opcodes {cmp, jcc, mov} → 100% Jaccard',
            font_size=12, color=RED, bold=True)

# ── Right: C source + assembly ──
add_code_box(sl, MID, Inches(1.5), Inches(3.5), Inches(1.2),
    'for (size_t i = 1; i+1 < n; i++) {\n'
    '    if (data[i] > data[i-1]\n'
    '     && data[i] > data[i+1]) {\n'
    '        uint64_t lmin = data[i];',
    font_size=10, title='C source — raw pointer, no checks', title_color=C_CLR)

# Show 2 diverse C blocks — pick blocks with different opcode sets to show diversity
c_blks = get_blocks(bc_c, 'O0')
_y_c = Inches(2.85)
add_textbox(sl, MID, _y_c, Inches(3.5), Inches(0.3),
            'C assembly — no bounds checks, blocks are diverse:', font_size=12, color=YELLOW, bold=True)
_y_c += Inches(0.3)
_c_shown = 0
_c_seen_ops = set()
for b in c_blks:
    ops_key = tuple(sorted(set(b.get('opcodes', []))))
    src = b.get('src_lines', [])
    src_label = f'line {src[0]}' if src else ''
    if 3 <= b['n_insns'] <= 8 and ops_key not in _c_seen_ops and _c_shown < 2:
        _c_seen_ops.add(ops_key)
        h = asm_box_height(b)
        add_code_box(sl, MID, _y_c, Inches(3.5), h,
            fmt_asm_full(b), font_size=9,
            title=f'{b["n_insns"]} insns ({src_label})',
            title_color=C_CLR)
        _y_c += h + Inches(0.05)
        _c_shown += 1

add_textbox(sl, MID, _y_c, Inches(3.5), Inches(0.5),
            '↑ No bounds checks — blocks have different opcodes',
            font_size=12, color=GREEN, bold=True)

# ── Results panel ──
# ── Results panel ──
add_rect(sl, Inches(10.0), Inches(1.5), RES_W, Inches(5.5), DARK_BG)
add_textbox(sl, Inches(10.1), Inches(1.6), RES_W - Inches(0.2), Inches(0.35),
            'bc_08 Accuracy', font_size=15, color=BLUE, bold=True)
_yr = 2.1
for m_name, m_label in [('value', 'Value'), ('opcodes', 'Opcodes'), ('constants', 'Constants')]:
    ra = bc_r['methods'][m_name]['accuracy']
    ca = bc_c['methods'][m_name]['accuracy']
    add_rect(sl, Inches(10.1), Inches(_yr), RES_W - Inches(0.2), Inches(0.4),
             RGBColor(0xEE, 0xF0, 0xF3),
             f'{m_label}: R {ra:.0%} | C {ca:.0%}', font_size=12, font_color=WHITE)
    _yr += 0.5
add_textbox(sl, Inches(10.1), Inches(_yr + 0.2), RES_W - Inches(0.2), Inches(2.5),
            '9 bounds check blocks:\n'
            'each does cmp index,len;\n'
            'jae panic. All identical\n'
            'opcodes {cmp, jcc, mov}.\n\n'
            'C: raw pointer, no checks →\n'
            'diverse opcodes → easy to match.',
            font_size=11, color=GRAY)

# ── Speaker notes for bc_08 ──
_bc08_notes = sl.notes_slide.notes_text_frame
_bc08_notes.text = (
    "=== FULL SOURCE CODE ===\n"
    "\n"
    "--- Rust: bc_08 (rust_crate/src/main.rs:963-980) ---\n"
    "fn bc_08(data: &[u64], min_prominence: u64) -> Vec<(usize, u64)> {\n"
    "    let n = data.len();\n"
    "    if n < 3 { return vec![]; }\n"
    "    let mut peaks = Vec::new();\n"
    "    for i in 1..n - 1 {\n"
    "        if data[i] > data[i - 1] && data[i] > data[i + 1] {\n"
    "            let mut left_min = data[i];\n"
    "            for j in (0..i).rev() {\n"
    "                if data[j] < left_min { left_min = data[j]; }\n"
    "                if data[j] > data[i] { break; }\n"
    "            }\n"
    "            if data[i] - left_min >= min_prominence {\n"
    "                peaks.push((i, data[i]));\n"
    "            }\n"
    "        }\n"
    "    }\n"
    "    peaks\n"
    "}\n"
    "\n"
    "--- C: bc_08 (c_src/bench.c:768-781) ---\n"
    "uint64_t bc_08(const uint64_t *data, size_t n, uint64_t min_prom) {\n"
    "    uint64_t count = 0;\n"
    "    for (size_t i = 1; i+1 < n; i++) {\n"
    "        if (data[i] > data[i-1] && data[i] > data[i+1]) {\n"
    "            uint64_t lmin = data[i];\n"
    "            for (size_t j = i; j > 0; j--) {\n"
    "                if (data[j-1] < lmin) lmin = data[j-1];\n"
    "                if (data[j-1] > data[i]) break;\n"
    "            }\n"
    "            if (data[i] - lmin >= min_prom) count++;\n"
    "        }\n"
    "    }\n"
    "    return count;\n"
    "}\n"
    "\n"
    "=== SLIDE ASSEMBLY BLOCKS — COMPLETE + SOURCE MAPPING ===\n"
    "\n"
    "Each bounds check block does: load index, load slice.len(), cmp, jae panic.\n"
    "If index < len: fall through to actual memory access (success).\n"
    "If index >= len: jump to panic handler (lea panic_msg + call panic).\n"
    "\n"
    "--- Rust BOUNDS_CHECK block 1 (for data[i] access on line 968) ---\n"
    "Source line 968: if data[i] > data[i - 1] && data[i] > data[i + 1]\n"
    "  rax = index (i), rcx = slice.len(). If i < len, continue; else panic.\n"
    "0x46f085: mov      rax, qword ptr [rsp + 0x68]\n"
    "0x46f08a: mov      rcx, qword ptr [rsp + 0xa8]\n"
    "0x46f092: cmp      rax, rcx\n"
    "0x46f095: jb       0x46f0ae\n"
    "  → jb (jump if below): index < len → success, proceed to access data[i]\n"
    "\n"
    "--- Rust BOUNDS_CHECK block 2 (for data[i-1] access on line 968) ---\n"
    "Source line 968: if data[i] > data[i - 1] && data[i] > data[i + 1]\n"
    "  rax = index (i-1 or i+1), rcx = slice.len(). If >= len, jump to panic.\n"
    "0x46f0ed: mov      rax, qword ptr [rsp + 0x70]\n"
    "0x46f0f2: mov      rcx, qword ptr [rsp + 0xa8]\n"
    "0x46f0fa: cmp      rax, rcx\n"
    "0x46f0fd: jae      0x46f128\n"
    "  → jae (jump if above/equal): index >= len → panic\n"
    "  Target 0x46f128 is: lea rdx,[rip+panic_msg]; call panic_handler\n"
    "\n"
    "--- Rust BOUNDS_CHECK block 3 (for data[i+1] access on line 968) ---\n"
    "Source line 968: same line, third array access\n"
    "0x46f14a: mov      rax, qword ptr [rsp + 0x58]\n"
    "0x46f14f: mov      rcx, qword ptr [rsp + 0xa8]\n"
    "0x46f157: cmp      rax, rcx\n"
    "0x46f15a: jb       0x46f173\n"
    "\n"
    "--- Remaining 6 Rust BOUNDS_CHECK blocks (not shown on slide) ---\n"
    "Block 4 [line 969: let mut left_min = data[i] — check i < len]:\n"
    "0x46f1b2: mov      rax, qword ptr [rsp + 0x70]\n"
    "0x46f1b7: mov      rcx, qword ptr [rsp + 0xa8]\n"
    "0x46f1bf: cmp      rax, rcx\n"
    "0x46f1c2: jae      0x46f1f2\n"
    "\n"
    "Block 5 [line 974: if data[i] - left_min — check i < len]:\n"
    "0x46f2ae: mov      rax, qword ptr [rsp + 0x70]\n"
    "0x46f2b3: mov      rcx, qword ptr [rsp + 0xa8]\n"
    "0x46f2bb: cmp      rax, rcx\n"
    "0x46f2be: jb       0x46f3fa\n"
    "\n"
    "Block 6 [line 972: if data[j] > data[i] — check j < len]:\n"
    "0x46f308: mov      rax, qword ptr [rsp + 0x18]\n"
    "0x46f30d: mov      rcx, qword ptr [rsp + 0xa8]\n"
    "0x46f315: cmp      rax, rcx\n"
    "0x46f318: jb       0x46f36b\n"
    "\n"
    "Block 7 [line 971: if data[j] < left_min — check j < len]:\n"
    "0x46f31c: mov      rax, qword ptr [rsp + 0x18]\n"
    "0x46f321: mov      rcx, qword ptr [rsp + 0xa8]\n"
    "0x46f329: cmp      rax, rcx\n"
    "0x46f32c: jae      0x46f349\n"
    "\n"
    "Block 8 [line 974: data[i] access — check i < len]:\n"
    "0x46f447: mov      rax, qword ptr [rsp + 8]\n"
    "0x46f44c: mov      rcx, qword ptr [rsp + 0xa0]\n"
    "0x46f454: cmp      rax, rcx\n"
    "0x46f457: jae      0x46f473\n"
    "\n"
    "Block 9 [line 975: peaks.push((i, data[i])) — check i < len]:\n"
    "0x46f473: mov      rax, qword ptr [rsp + 0x70]\n"
    "0x46f478: mov      rcx, qword ptr [rsp + 0xa8]\n"
    "0x46f480: cmp      rax, rcx\n"
    "0x46f483: jae      0x46f4a5\n"
    "\n"
    "--- Example panic target (where jae jumps to on bounds check failure) ---\n"
    "Target of block 2's jae 0x46f128:\n"
    "0x46f128: mov      rsi, qword ptr [rsp + 0xa8]\n"
    "0x46f130: mov      rdi, qword ptr [rsp + 0x70]\n"
    "0x46f135: lea      rdx, [rip + 0x58cf4]\n"
    "0x46f13c: mov      rax, qword ptr [rip + 0x5d43d]\n"
    "0x46f143: call     rax\n"
    "  → loads panic message (\"index out of bounds\") and calls core::panicking::panic_bounds_check\n"
    "\n"
    "--- C block 1 shown (8 insns) ---\n"
    "Source lines 768-770: function prologue + variable init + for loop start\n"
    "  C has NO array bounds checking — raw pointer access.\n"
    "0x405453: push     rbp\n"
    "0x405454: mov      rbp, rsp\n"
    "0x405457: mov      qword ptr [rbp - 0x28], rdi\n"
    "0x40545b: mov      qword ptr [rbp - 0x30], rsi\n"
    "0x40545f: mov      qword ptr [rbp - 0x38], rdx\n"
    "0x405463: mov      qword ptr [rbp - 8], 0\n"
    "0x40546b: mov      qword ptr [rbp - 0x10], 1\n"
    "0x405473: jmp      0x4055aa\n"
    "\n"
    "--- C block 2 shown (8 insns) ---\n"
    "Source line 774: if (data[j-1] < lmin) lmin = data[j-1];\n"
    "  Direct pointer arithmetic: shl+lea to compute offset, then load.\n"
    "  No cmp-index-vs-length — just raw memory access.\n"
    "0x40550a: mov      rax, qword ptr [rbp - 0x20]\n"
    "0x40550e: shl      rax, 3\n"
    "0x405512: lea      rdx, [rax - 8]\n"
    "0x405516: mov      rax, qword ptr [rbp - 0x28]\n"
    "0x40551a: add      rax, rdx\n"
    "0x40551d: mov      rax, qword ptr [rax]\n"
    "0x405520: cmp      qword ptr [rbp - 0x18], rax\n"
    "0x405524: jbe      0x405540"
)

# ═══════════════════════════════════════════════════════════════════════════
# Slide 7: Example 2 — ? Operator (qm_06)
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
slide_title(sl, 'Example: ? Operator (qm_06)',
            'Each ? generates identical error-propagation blocks — C has no equivalent')

qm_r = get_fn('qm_06', 'Rust')
qm_c = get_fn('qm_06', 'C')

# ── Left: Rust source + assembly ──
add_code_box(sl, LEFT, Inches(1.5), BLK_W, Inches(1.3),
    'fn qm_06(data: &[u64]) -> Option<u64> {\n'
    '    let first = data.first()?;          // ← ? generates 2 blocks\n'
    '    let last = data.last()?;            // ← ? generates 2 blocks\n'
    '    let mid = data.get(mid_idx)?;       // ← ? generates 2 blocks\n'
    '    let range = last.checked_sub(*first)?;\n'
    '    // ... 6 uses of ? total',
    font_size=10, title='Rust — each ? inserts: discriminant check + error return', title_color=RUST_CLR)

qm_r_blks = get_blocks(qm_r, 'O0')

# Find the identical [jmp, mov] error-propagation blocks (they all jump to same addr)
_qm_err_blocks = []
for b in qm_r_blks:
    ops = tuple(sorted(set(b.get('opcodes', []))))
    asm_text = ' '.join(b.get('asm', []))
    if ops == ('jmp', 'mov') and b['n_insns'] <= 4:
        _qm_err_blocks.append(b)

# Find a discriminant-check block (cmp/test + jcc)
_qm_check = None
for b in qm_r_blks:
    ops = set(b.get('opcodes', []))
    if ('cmp' in ops or 'test' in ops) and ('je' in b.get('opcodes', []) or 'jcc' in ops) and b['n_insns'] <= 10:
        _qm_check = b
        break

_y = Inches(2.95)
add_textbox(sl, LEFT, _y, BLK_W, Inches(0.3),
            f'Rust assembly — {len(_qm_err_blocks)} identical error-return blocks (of {len(qm_r_blks)} O0 blocks):', font_size=12, color=YELLOW, bold=True)
_y += Inches(0.3)

# Show the check block first
if _qm_check:
    h = asm_box_height(_qm_check)
    add_code_box(sl, LEFT, _y, Inches(5.0), h,
        fmt_asm_full(_qm_check), font_size=9,
        title='discriminant check: is it None?',
        title_color=PURPLE)
    _y += h + Inches(0.05)

# Show 3 of the identical error-propagation blocks
for i, blk in enumerate(_qm_err_blocks[:3]):
    h = asm_box_height(blk)
    label = f'error-return block {i+1} — same opcodes {{jmp, mov}}'
    add_code_box(sl, LEFT, _y, Inches(5.0), h,
        fmt_asm_full(blk), font_size=9,
        title=label,
        title_color=YELLOW)
    _y += h + Inches(0.05)

add_textbox(sl, LEFT, _y, BLK_W, Inches(0.5),
            f'↑ {len(_qm_err_blocks)} blocks with identical opcodes {{jmp, mov}} — matcher pairs wrong ones',
            font_size=12, color=RED, bold=True)

# ── Right: C source + assembly ──
add_code_box(sl, MID, Inches(1.5), Inches(3.5), Inches(1.4),
    'int qm_06(const uint64_t *data, size_t n,\n'
    '           uint64_t *out) {\n'
    '    if (n < 3) return 0; // ← just return 0\n'
    '    uint64_t first=data[0], last=data[n-1];\n'
    '    if (last < first) return 0;\n'
    '    uint64_t range = last - first;\n'
    '    // No Option, no ? — direct access + return code',
    font_size=10, title='C — return code instead of Option/?', title_color=C_CLR)

c_blks = get_blocks(qm_c, 'O0')
_y_c = Inches(2.65)
add_textbox(sl, MID, _y_c, Inches(3.5), Inches(0.3),
            f'C assembly — {len(c_blks)} O0 blocks, diverse opcodes:', font_size=12, color=YELLOW, bold=True)
_y_c += Inches(0.3)
_c_shown = 0
_c_seen_ops2 = set()
for b in c_blks:
    ops_key = tuple(sorted(set(b.get('opcodes', []))))
    src = b.get('src_lines', [])
    src_label = f'line {src[0]}' if src else ''
    if 2 <= b['n_insns'] <= 8 and ops_key not in _c_seen_ops2 and _c_shown < 3:
        _c_seen_ops2.add(ops_key)
        h = asm_box_height(b)
        add_code_box(sl, MID, _y_c, Inches(3.5), h,
            fmt_asm_full(b), font_size=9,
            title=f'{b["n_insns"]} insns ({src_label})',
            title_color=C_CLR)
        _y_c += h + Inches(0.05)
        _c_shown += 1

add_textbox(sl, MID, _y_c, Inches(3.5), Inches(0.5),
            '↑ No error-return blocks → diverse opcodes → easy to match',
            font_size=12, color=GREEN, bold=True)

# ── Results panel ──
add_rect(sl, Inches(10.0), Inches(1.5), RES_W, Inches(5.5), DARK_BG)
add_textbox(sl, Inches(10.1), Inches(1.6), RES_W - Inches(0.2), Inches(0.35),
            'qm_06 Accuracy', font_size=15, color=BLUE, bold=True)
_yr = 2.1
for m_name, m_label in [('value', 'Value'), ('opcodes', 'Opcodes'), ('constants', 'Constants')]:
    ra = qm_r['methods'][m_name]['accuracy']
    ca = qm_c['methods'][m_name]['accuracy']
    add_rect(sl, Inches(10.1), Inches(_yr), RES_W - Inches(0.2), Inches(0.4),
             RGBColor(0xEE, 0xF0, 0xF3),
             f'{m_label}: R {ra:.0%} | C {ca:.0%}', font_size=12, font_color=WHITE)
    _yr += 0.5
add_textbox(sl, Inches(10.1), Inches(_yr + 0.2), RES_W - Inches(0.2), Inches(2.5),
            '49 Rust O0 blocks, only 22%\n'
            'unique opcode sets.\n\n'
            '6 uses of ? → 8 identical\n'
            'error-return blocks.\n\n'
            'C: 9 O0 blocks, 89% unique.\n'
            'No ? operator → diverse blocks.',
            font_size=11, color=GRAY)

# ── Speaker notes for qm_06 ──
_qm06_notes = sl.notes_slide.notes_text_frame
_qm06_notes.text = (
    "=== FULL SOURCE CODE ===\n"
    "\n"
    "--- Rust: qm_06 (rust_crate/src/main.rs:1280-1293) ---\n"
    "fn qm_06(data: &[u64]) -> Option<u64> {\n"
    "    let first = data.first()?;\n"
    "    let last = data.last()?;\n"
    "    if data.len() < 3 { return None; }\n"
    "    let mid_idx = data.len() / 2;\n"
    "    let mid = data.get(mid_idx)?;\n"
    "    let range = last.checked_sub(*first)?;\n"
    "    let deviation = if *mid > *first + range / 2 {\n"
    "        mid.checked_sub(*first + range / 2)?\n"
    "    } else {\n"
    "        (*first + range / 2).checked_sub(*mid)?\n"
    "    };\n"
    "    Some(range.wrapping_add(deviation))\n"
    "}\n"
    "\n"
    "--- C: qm_06 (c_src/bench.c:1025-1033) ---\n"
    "int qm_06(const uint64_t *data, size_t n, uint64_t *out) {\n"
    "    if (n < 3) return 0;\n"
    "    uint64_t first = data[0], last = data[n-1], mid = data[n/2];\n"
    "    if (last < first) return 0;\n"
    "    uint64_t range = last - first;\n"
    "    uint64_t center = first + range / 2;\n"
    "    uint64_t dev = mid > center ? mid - center : center - mid;\n"
    "    *out = range + dev; return 1;\n"
    "}\n"
    "\n"
    "=== SLIDE ASSEMBLY BLOCKS — COMPLETE + SOURCE MAPPING ===\n"
    "\n"
    "--- Rust: discriminant check block (shown on slide) ---\n"
    "Source: compiler-generated for data.first()? (line 1281)\n"
    "0x47f06e: mov      qword ptr [rsp + 0x68], rax\n"
    "0x47f073: mov      rdx, qword ptr [rsp + 0x68]\n"
    "0x47f078: xor      eax, eax\n"
    "0x47f07a: mov      ecx, 1\n"
    "0x47f07f: cmp      rdx, 0\n"
    "0x47f083: cmove    rax, rcx\n"
    "0x47f087: test     rax, 1\n"
    "0x47f08d: je       0x47f0a3\n"
    "\n"
    "--- Rust: error-return block 1 (shown on slide) ---\n"
    "Source: compiler-generated for data.first()? (line 1281) error path\n"
    "0x47f094: mov      qword ptr [rsp + 0x58], rax\n"
    "0x47f099: mov      qword ptr [rsp + 0x60], rdx\n"
    "0x47f09e: jmp      0x47f33f\n"
    "\n"
    "--- Rust: error-return block 2 (shown on slide) ---\n"
    "Source: compiler-generated for data.last()? (line 1282) error path\n"
    "0x47f0f2: mov      qword ptr [rsp + 0x58], rax\n"
    "0x47f0f7: mov      qword ptr [rsp + 0x60], rdx\n"
    "0x47f0fc: jmp      0x47f33f\n"
    "\n"
    "--- Rust: error-return block 3 (shown on slide) ---\n"
    "Source line 1283: if data.len() < 3 { return None; }\n"
    "0x47f166: mov      qword ptr [rsp + 0x58], 0\n"
    "0x47f16f: jmp      0x47f33f\n"
    "\n"
    "--- Remaining 5 Rust error-return blocks (not shown on slide) ---\n"
    "Block 4 [compiler-generated for data.get(mid_idx)? line 1285]:\n"
    "0x47f179: mov      qword ptr [rsp + 0x58], rax\n"
    "0x47f17e: mov      qword ptr [rsp + 0x60], rdx\n"
    "0x47f183: jmp      0x47f33f\n"
    "\n"
    "Block 5 [compiler-generated for last.checked_sub(*first)? line 1286]:\n"
    "0x47f1dd: mov      qword ptr [rsp + 0x58], rax\n"
    "0x47f1e2: mov      qword ptr [rsp + 0x60], rdx\n"
    "0x47f1e7: jmp      0x47f33f\n"
    "\n"
    "Block 6 [compiler-generated for mid.checked_sub()? line 1288]:\n"
    "0x47f2ed: mov      qword ptr [rsp + 0x58], rax\n"
    "0x47f2f2: mov      qword ptr [rsp + 0x60], rdx\n"
    "0x47f2f7: jmp      0x47f33f\n"
    "\n"
    "Block 7 [compiler-generated for .checked_sub(*mid)? line 1290]:\n"
    "0x47f39c: mov      qword ptr [rsp + 0x58], rax\n"
    "0x47f3a1: mov      qword ptr [rsp + 0x60], rdx\n"
    "0x47f3a6: jmp      0x47f33f\n"
    "\n"
    "Block 8 [lines 1287-1288: deviation branch]:\n"
    "0x47f3a8: mov      rax, qword ptr [rsp + 0xa0]\n"
    "0x47f3b0: mov      qword ptr [rsp + 0x100], rax\n"
    "0x47f3b8: mov      qword ptr [rsp + 0x90], rax\n"
    "0x47f3c0: jmp      0x47f311\n"
    "\n"
    "--- C block 1 shown (BODY, 7 insns) ---\n"
    "Source lines 1025-1026: int qm_06(...) { if (n < 3) return 0;\n"
    "0x406bfd: push     rbp\n"
    "0x406bfe: mov      rbp, rsp\n"
    "0x406c01: mov      qword ptr [rbp - 0x38], rdi\n"
    "0x406c05: mov      qword ptr [rbp - 0x40], rsi\n"
    "0x406c09: mov      qword ptr [rbp - 0x48], rdx\n"
    "0x406c0d: cmp      qword ptr [rbp - 0x40], 2\n"
    "0x406c12: ja       0x406c1e\n"
    "\n"
    "--- C block 2 shown (BODY, 2 insns) ---\n"
    "Source line 1026: return 0 (when n < 3)\n"
    "0x406c14: mov      eax, 0\n"
    "0x406c19: jmp      0x406cc9\n"
    "\n"
    "--- C block 3 shown (BODY, 2 insns) ---\n"
    "Source line 1028: return 0 (when last < first)\n"
    "0x406c6a: mov      eax, 0\n"
    "0x406c6f: jmp      0x406cc9"
)

# ═══════════════════════════════════════════════════════════════════════════
# Slide 8: Example 3 — Panic Paths (pu_08)
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
slide_title(sl, 'Example: Panic Paths (pu_08)',
            '.unwrap() inserts check + panic blocks — panic blocks look like normal call blocks')

pu_r = get_fn('pu_08', 'Rust')
pu_c = get_fn('pu_08', 'C')

# ── Left: Rust source + assembly ──
add_code_box(sl, LEFT, Inches(1.5), BLK_W, Inches(1.5),
    'fn pu_08(s: &str) -> u64 {\n'
    '    let parts: Vec<&str> = s.split_whitespace().collect();\n'
    '    let first = parts.first().unwrap(); // check+panic\n'
    '    let last = parts.last().unwrap();   // check+panic\n'
    '    let mid = parts.get(parts.len()/2).unwrap(); // check+panic\n'
    '// Each .unwrap(): check discriminant → if None, call panic\n'
    '// Panic block uses {lea, mov, call} — same as normal calls',
    font_size=10, title='Rust source — .unwrap() inserts check + panic blocks', title_color=RUST_CLR)

_pu_blks = get_blocks(pu_r, 'O0')

# Find the check block (at 0x47bf85)
_pu_check = None
for b in _pu_blks:
    if b['addr'] == '0x47bf85':
        _pu_check = b
        break

# Find the panic block (at 0x47bfb4)
_pu_panic = None
for b in _pu_blks:
    if b['addr'] == '0x47bfb4':
        _pu_panic = b
        break

# Find a normal function call block for contrast (e.g. collect() call)
_pu_normal = None
for b in _pu_blks:
    src = b.get('src_lines', [])
    asm_text = ' '.join(b.get('asm', []))
    if 'call' in asm_text and 'rip' not in asm_text and 'lea' in asm_text and 3 <= b['n_insns'] <= 5:
        _pu_normal = b
        break

_y = Inches(3.15)
add_textbox(sl, LEFT, _y, BLK_W, Inches(0.3),
            'Rust assembly — .unwrap() expansion + a normal call block:', font_size=12, color=YELLOW, bold=True)
_y += Inches(0.3)

_unwrap_parts = [
    ('check: is Option None? (cmp rdx,0 → if None, fall to panic)', _pu_check, PURPLE),
    ('panic block: lea panic_msg + call panic — opcodes: {call, lea, mov}', _pu_panic, RED),
    ('normal call block (collect) — opcodes: {call, lea, mov} ← SAME!', _pu_normal, C_CLR),
]
for label, blk, clr in _unwrap_parts:
    if blk:
        h = asm_box_height(blk)
        add_code_box(sl, LEFT, _y, Inches(5.0), h,
            fmt_asm_full(blk), font_size=9,
            title=label, title_color=clr)
        _y += h + Inches(0.05)

add_textbox(sl, LEFT, _y, BLK_W, Inches(0.5),
            '↑ Panic block and normal call block: same opcodes {call, lea, mov} — matcher can\'t distinguish',
            font_size=12, color=RED, bold=True)

# ── Right: C source + assembly ──
add_code_box(sl, MID, Inches(1.5), Inches(3.5), Inches(1.5),
    'uint64_t pu_08(const char *s) {\n'
    '    char *words[256]; // strtok tokenize\n'
    '    // Direct access — no .unwrap(), no panic:\n'
    '    for (i=0; words[0][i]; i++) sum += words[0][i];\n'
    '    for (i=0; words[nw/2][i]; i++) sum += words[nw/2][i];\n'
    '    for (i=0; words[nw-1][i]; i++) sum += words[nw-1][i];\n'
    '    return sum; }',
    font_size=10, title='C — direct array access, no panic paths', title_color=C_CLR)

c_blks = get_blocks(pu_c, 'O0')
_y_c = Inches(2.85)
add_textbox(sl, MID, _y_c, Inches(3.5), Inches(0.3),
            'C assembly — no panic paths, blocks are diverse:', font_size=12, color=YELLOW, bold=True)
_y_c += Inches(0.3)
_c_shown = 0
_c_seen_ops3 = set()
for b in c_blks:
    ops_key = tuple(sorted(set(b.get('opcodes', []))))
    src = b.get('src_lines', [])
    src_label = f'line {src[0]}' if src else ''
    if 2 <= b['n_insns'] <= 6 and ops_key not in _c_seen_ops3 and _c_shown < 2:
        _c_seen_ops3.add(ops_key)
        h = asm_box_height(b)
        add_code_box(sl, MID, _y_c, Inches(3.5), h,
            fmt_asm_full(b), font_size=9,
            title=f'{b["n_insns"]} insns ({src_label})',
            title_color=C_CLR)
        _y_c += h + Inches(0.05)
        _c_shown += 1

add_textbox(sl, MID, _y_c, Inches(3.5), Inches(0.5),
            '↑ No panic paths → diverse opcodes → easy to match',
            font_size=12, color=GREEN, bold=True)

# Bottom stat
add_textbox(sl, LEFT, Inches(6.7), Inches(9.5), Inches(0.4),
            'pu_08 Rust O0: 43 blocks, only 10 unique opcode sets. 11× {call, lea, mov}, 12× {jmp, mov}. Matcher can\'t tell them apart.',
            font_size=13, color=WHITE)

# ── Results panel ──
add_rect(sl, Inches(10.0), Inches(1.5), RES_W, Inches(5.5), DARK_BG)
add_textbox(sl, Inches(10.1), Inches(1.6), RES_W - Inches(0.2), Inches(0.35),
            'pu_08 Accuracy', font_size=15, color=BLUE, bold=True)
_yr = 2.1
for m_name, m_label in [('value', 'Value'), ('opcodes', 'Opcodes'), ('constants', 'Constants')]:
    ra = pu_r['methods'][m_name]['accuracy']
    ca = pu_c['methods'][m_name]['accuracy']
    add_rect(sl, Inches(10.1), Inches(_yr), RES_W - Inches(0.2), Inches(0.4),
             RGBColor(0xEE, 0xF0, 0xF3),
             f'{m_label}: R {ra:.0%} | C {ca:.0%}', font_size=12, font_color=WHITE)
    _yr += 0.5
add_textbox(sl, Inches(10.1), Inches(_yr + 0.2), RES_W - Inches(0.2), Inches(2.5),
            '3× .unwrap() = 3 panic blocks\n'
            'each with {call, lea, mov}.\n\n'
            '11 blocks share {call, lea, mov}\n'
            '(panic + normal calls mixed).\n\n'
            'C: no panic paths →\n'
            'blocks are diverse.',
            font_size=11, color=GRAY)

# ── Speaker notes for pu_08 ──
_pu08_notes = sl.notes_slide.notes_text_frame
_pu08_notes.text = (
    "=== FULL SOURCE CODE ===\n"
    "\n"
    "--- Rust: pu_08 (rust_crate/src/main.rs:1601-1611) ---\n"
    "fn pu_08(s: &str) -> u64 {\n"
    "    let parts: Vec<&str> = s.split_whitespace().collect();\n"
    "    let first = parts.first().unwrap();\n"
    "    let last = parts.last().unwrap();\n"
    "    let mid = parts.get(parts.len() / 2).unwrap();\n"
    "    let mut sum = 0u64;\n"
    "    for ch in first.chars().chain(mid.chars()).chain(last.chars()) {\n"
    "        sum = sum.wrapping_add(ch as u64);\n"
    "    }\n"
    "    sum\n"
    "}\n"
    "\n"
    "--- C: pu_08 (c_src/bench.c:1280-1289) ---\n"
    "uint64_t pu_08(const char *s) {\n"
    "    char *words[256]; size_t nw = 0;\n"
    "    char copy[1024]; strncpy(copy, s, 1023); copy[1023] = '\\0';\n"
    "    char *tok = strtok(copy, \" \");\n"
    "    while (tok && nw < 256) { words[nw++] = tok; tok = strtok(NULL, \" \"); }\n"
    "    uint64_t sum = 0;\n"
    "    for (size_t i = 0; words[0][i]; i++) sum += words[0][i];\n"
    "    for (size_t i = 0; words[nw/2][i]; i++) sum += words[nw/2][i];\n"
    "    for (size_t i = 0; words[nw-1][i]; i++) sum += words[nw-1][i];\n"
    "    return sum;\n"
    "}\n"
    "\n"
    "=== SLIDE ASSEMBLY BLOCKS — COMPLETE + SOURCE MAPPING ===\n"
    "\n"
    "--- Rust: check block (0x47bf85) — is Option None? (shown on slide) ---\n"
    "Source: compiler-generated for parts.first().unwrap() (line 1603)\n"
    "  cmp rdx,0: checks if Option is None. If not None → jne to 0x47bfc8 (success).\n"
    "  If None → falls through to panic block at 0x47bfb4.\n"
    "0x47bf85: mov      rax, qword ptr [rsp + 0xc0]\n"
    "0x47bf8d: mov      qword ptr [rsp + 0x288], rax\n"
    "0x47bf95: mov      rdx, qword ptr [rsp + 0x288]\n"
    "0x47bf9d: mov      eax, 1\n"
    "0x47bfa2: xor      ecx, ecx\n"
    "0x47bfa4: cmp      rdx, 0\n"
    "0x47bfa8: cmove    rax, rcx\n"
    "0x47bfac: test     rax, 1\n"
    "0x47bfb2: jne      0x47bfc8\n"
    "\n"
    "--- Rust: panic block (0x47bfb4) — lea panic_msg + call panic (shown on slide) ---\n"
    "Source: compiler-generated for .unwrap() panic path (line 1603)\n"
    "  Opcodes: {call, lea, mov} — SAME as normal function calls.\n"
    "  This is the key problem: matcher can't distinguish panic from real calls.\n"
    "0x47bfb4: lea      rdi, [rip + 0x4d39d]\n"
    "0x47bfbb: mov      rax, qword ptr [rip + 0x506fe]\n"
    "0x47bfc2: call     rax\n"
    "\n"
    "--- Rust: normal call block — collect() (shown on slide for CONTRAST) ---\n"
    "Source line 1602: let parts: Vec<&str> = s.split_whitespace().collect();\n"
    "  Opcodes: {call, lea, mov} — SAME as the panic block above!\n"
    "  This demonstrates why the matcher pairs them incorrectly.\n"
    "0x47bf02: mov      rsi, qword ptr [rsp + 0xc8]\n"
    "0x47bf0a: lea      rdi, [rsp + 0xe8]\n"
    "0x47bf12: mov      qword ptr [rsp + 0xd0], rdi\n"
    "0x47bf1a: call     0x42efa0\n"
    "\n"
    "--- C block 1 shown ---\n"
    "Source lines 1282-1283: strncpy(copy, s, 1023); copy[1023] = '\\0';\n"
    "0x407fd6: mov      byte ptr [rbp - 0x831], 0\n"
    "0x407fdd: lea      rax, [rbp - 0xc30]\n"
    "0x407fe4: mov      esi, 0x40b019\n"
    "0x407fe9: mov      rdi, rax\n"
    "0x407fec: call     0x401110\n"
    "\n"
    "--- C block 2 shown ---\n"
    "Source line 1283: strtok loop — store result and continue\n"
    "0x407ff1: mov      qword ptr [rbp - 0x10], rax\n"
    "0x407ff5: jmp      0x408022"
)

# ═══════════════════════════════════════════════════════════════════════════
# Slide 9: Overall Results — the big table
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
slide_title(sl, 'Results: C Wins Across All Features and Methods',
            'C wins 288 / 300 individual comparisons (100 pairs × 3 methods)')

# Main table header
headers = ['Feature', 'Rust\nBlocks', 'C\nBlocks', 'Value\nR / C', 'Opcodes\nR / C', 'Constants\nR / C']
col_widths = [Inches(2.8), Inches(1.3), Inches(1.3), Inches(2.2), Inches(2.2), Inches(2.2)]
col_x = [Inches(0.5)]
for w in col_widths[:-1]:
    col_x.append(col_x[-1] + w)

y = Inches(2.0)
row_h = Inches(0.4)

# Header row
for i, (hdr, w) in enumerate(zip(headers, col_widths)):
    add_rect(sl, col_x[i], y, w - Pt(2), row_h, RGBColor(0xEE, 0xF0, 0xF3),
             hdr, font_size=12, font_color=BLUE)

# Data rows
feat_data = [
    ('bc', 'Bounds Checking'),
    ('om', 'Safe Abstractions'),
    ('dg', 'Drop Glue / RAII'),
    ('qm', '? Operator'),
    ('pu', 'Panic Paths'),
]
y += row_h + Pt(2)
for prefix, feat_name in feat_data:
    r_fns = [f for f in DATA['functions'] if f['lang'] == 'Rust' and f['feature_prefix'] == prefix]
    c_fns_feat = [f for f in DATA['functions'] if f['lang'] == 'C' and f['feature_prefix'] == prefix]
    r_o0 = sum(f['n_o0'] for f in r_fns) / len(r_fns)
    c_o0 = sum(f['n_o0'] for f in c_fns_feat) / len(c_fns_feat)

    add_rect(sl, col_x[0], y, col_widths[0] - Pt(2), row_h, DARK_BG,
             feat_name, font_size=13, font_color=WHITE)
    add_rect(sl, col_x[1], y, col_widths[1] - Pt(2), row_h, DARK_BG,
             f'{r_o0:.0f}', font_size=13, font_color=RUST_CLR)
    add_rect(sl, col_x[2], y, col_widths[2] - Pt(2), row_h, DARK_BG,
             f'{c_o0:.0f}', font_size=13, font_color=C_CLR)

    for j, m in enumerate(['value', 'opcodes', 'constants']):
        r_acc = sum(f['methods'][m]['accuracy'] for f in r_fns) / len(r_fns)
        c_acc = sum(f['methods'][m]['accuracy'] for f in c_fns_feat) / len(c_fns_feat)
        add_rect(sl, col_x[3 + j], y, col_widths[3 + j] - Pt(2), row_h, DARK_BG,
                 f'{r_acc:.0%}  /  {c_acc:.0%}', font_size=13, font_color=WHITE)
    y += row_h + Pt(2)

# Average row
y += Pt(4)
add_rect(sl, col_x[0], y, col_widths[0] - Pt(2), row_h, RGBColor(0xEE, 0xF0, 0xF3),
         'AVERAGE', font_size=13, font_color=YELLOW)
add_rect(sl, col_x[1], y, col_widths[1] - Pt(2), row_h, RGBColor(0xEE, 0xF0, 0xF3),
         f'{r_avg_o0:.0f}', font_size=13, font_color=RUST_CLR)
add_rect(sl, col_x[2], y, col_widths[2] - Pt(2), row_h, RGBColor(0xEE, 0xF0, 0xF3),
         f'{c_avg_o0:.0f}', font_size=13, font_color=C_CLR)
for j, m in enumerate(['value', 'opcodes', 'constants']):
    r_acc = sum(f['methods'][m]['accuracy'] for f in rust_fns) / len(rust_fns)
    c_acc = sum(f['methods'][m]['accuracy'] for f in c_fns) / len(c_fns)
    add_rect(sl, col_x[3 + j], y, col_widths[3 + j] - Pt(2), row_h, RGBColor(0xEE, 0xF0, 0xF3),
             f'{r_acc:.0%}  /  {c_acc:.0%}', font_size=13, font_color=YELLOW)

# Bottom insights
y_bottom = Inches(5.2)
add_bullet_frame(sl, Inches(0.6), y_bottom, Inches(12), Inches(2.0), [
    ('Key patterns:', BLUE),
    ('• All Rust features hurt matching — safety blocks (bounds checks, error returns, panic paths) share the same opcodes', WHITE),
    ('• The ? operator and drop glue create the most identical blocks — worst Rust accuracy (~5-7%)', WHITE),
    ('• Opcodes is the best method overall — but still struggles when many blocks share the same instruction types', WHITE),
    ('• C has no equivalent safety overhead → blocks are inherently more diverse → easier to match', WHITE),
], font_size=14, spacing=Pt(3))


# ═══════════════════════════════════════════════════════════════════════════
# Slide 10: Findings & Takeaways
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
slide_title(sl, 'Findings')

findings = [
    ('1. Rust binaries are significantly harder to match than C',
     'C wins 288/300 comparisons. Gap ranges from 3× to 10× depending on feature and method.',
     RUST_CLR),
    ('2. Root cause: safety blocks are too similar to each other and to computation blocks',
     'Bounds checks, drop glue, panic paths all use generic opcodes (cmp, call, lea, mov) — the same\n'
     'instructions used in real computation. High frequency + low diversity → wrong matches.',
     YELLOW),
    ('3. Drop glue is the hardest, panic paths the easiest (among Rust features)',
     'Drop blocks (lea+call) are indistinguishable from normal function calls.\n'
     'Panic blocks at least survive O2, providing some stable anchor points.',
     PURPLE),
    ('4. Opcode Jaccard is the most robust feature, but still limited',
     'When many blocks share the same opcode set, even the best feature can\'t tell them apart.',
     GREEN),
    ('5. Possible direction: Rust-aware pre-filtering',
     'Identify safety blocks before matching (by call target, branch pattern) and handle them separately.',
     BLUE),
]

y = 1.6
for title, body, color in findings:
    add_textbox(sl, Inches(0.6), Inches(y), Inches(12), Inches(0.35),
                title, font_size=17, color=color, bold=True)
    add_textbox(sl, Inches(0.9), Inches(y + 0.35), Inches(11.5), Inches(0.7),
                body, font_size=14, color=WHITE)
    y += 1.1


# ═══════════════════════════════════════════════════════════════════════════
# Slide 11: Discussion / What's Next
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
slide_title(sl, 'Discussion Topics')

add_bullet_frame(sl, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.5), [
    ('Possible improvements:', BLUE),
    ('• Pre-identify safety blocks (bounds check, drop, panic)\n  and filter or weight them differently', WHITE),
    ('• Hybrid matching: combine structural + value features\n  with language-aware weights', WHITE),
    ('• Test on real-world Rust projects\n  (coreutils, ripgrep) beyond synthetic corpus', WHITE),
    ('• Try other features: CFG structure similarity,\n  string references, callee names', WHITE),
], font_size=16, spacing=Pt(8))

add_bullet_frame(sl, Inches(7.0), Inches(1.8), Inches(5.5), Inches(4.5), [
    ('Limitations of this study:', BLUE),
    ('• Synthetic corpus — 100 hand-written function pairs\n  may not represent all real-world patterns', WHITE),
    ('• Single compiler (rustc 1.92 + gcc 11.4)\n  — results may vary across toolchains', WHITE),
    ('• Only O0 vs O2 — O1, Os, Oz not tested', WHITE),
    ('• Block-level matching only —\n  function-level matching not evaluated here', WHITE),
], font_size=16, spacing=Pt(8))



# ═══════════════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════════════
out_path = Path('experiments/rust_features/results/slides.pptx')
prs.save(str(out_path))
print(f"Saved: {out_path} ({out_path.stat().st_size / 1024:.0f} KB)")
